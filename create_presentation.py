"""
CareVance HRMS Comprehensive Presentation Generator
Creates a detailed PowerPoint presentation covering all modules and features
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_AUTO_SIZE

# Color scheme - Professional Blue/Green theme
PRIMARY_COLOR = RGBColor(0, 150, 136)  # Teal
SECONDARY_COLOR = RGBColor(63, 81, 181)  # Indigo
ACCENT_COLOR = RGBColor(255, 152, 0)  # Orange
DARK_COLOR = RGBColor(33, 37, 41)  # Dark Gray
LIGHT_BG = RGBColor(240, 248, 255)  # Light Blue
WHITE = RGBColor(255, 255, 255)

def add_title_slide(prs, title, subtitle=""):
    """Add a title slide"""
    slide_layout = prs.slide_layouts[0]  # Title slide
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Set subtitle
    if subtitle:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle
        subtitle_frame = subtitle_shape.text_frame
        subtitle_frame.paragraphs[0].font.size = Pt(24)
        subtitle_frame.paragraphs[0].font.color.rgb = DARK_COLOR
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, section_title):
    """Add a section divider slide"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add background shape
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5)
    )
    background.fill.solid()
    background.fill.fore_color.rgb = PRIMARY_COLOR
    background.line.fill.background()
    
    # Add section title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.333), Inches(1.5))
    tf = title_box.text_frame
    tf.text = section_title
    p = tf.paragraphs[0]
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = WHITE
    p.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_items, subtitle=""):
    """Add a content slide with bullet points"""
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    
    # Set title
    title_shape = slide.shapes.title
    title_shape.text = title
    title_frame = title_shape.text_frame
    title_frame.paragraphs[0].font.size = Pt(36)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].font.color.rgb = PRIMARY_COLOR
    
    # Set subtitle if provided
    if subtitle:
        title_frame.add_paragraph()
        title_frame.paragraphs[1].text = subtitle
        title_frame.paragraphs[1].font.size = Pt(18)
        title_frame.paragraphs[1].font.color.rgb = DARK_COLOR
    
    # Set content
    body_shape = slide.placeholders[1]
    tf = body_shape.text_frame
    tf.clear()
    
    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = item
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_COLOR
        p.level = 0 if not item.startswith("  ") else 1
        p.space_after = Pt(12)
    
    return slide

def add_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Add a slide with two columns"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.9), Inches(5.5))
    tf_left = left_box.text_frame
    tf_left.text = left_title
    p = tf_left.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    for item in left_items:
        p = tf_left.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(8)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.9), Inches(1.5), Inches(5.9), Inches(5.5))
    tf_right = right_box.text_frame
    tf_right.text = right_title
    p = tf_right.paragraphs[0]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = SECONDARY_COLOR
    
    for item in right_items:
        p = tf_right.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_COLOR
        p.space_after = Pt(8)
    
    return slide

def add_feature_grid_slide(prs, title, features):
    """Add a slide with features in a grid layout"""
    slide_layout = prs.slide_layouts[5]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Add title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.333), Inches(1))
    tf = title_box.text_frame
    tf.text = title
    p = tf.paragraphs[0]
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = PRIMARY_COLOR
    p.alignment = PP_ALIGN.CENTER
    
    # Calculate grid layout
    num_features = len(features)
    cols = 3
    rows = (num_features + cols - 1) // cols
    
    col_width = 4
    row_height = 1.2
    start_x = 0.5
    start_y = 1.5
    
    for i, feature in enumerate(features):
        col = i % cols
        row = i // cols
        
        x = start_x + col * col_width
        y = start_y + row * row_height
        
        # Add feature box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(y),
            Inches(3.8), Inches(1.1)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = LIGHT_BG
        box.line.color.rgb = PRIMARY_COLOR
        box.line.width = Pt(2)
        
        # Add feature text
        text_box = slide.shapes.add_textbox(
            Inches(x + 0.1), Inches(y + 0.1),
            Inches(3.6), Inches(0.9)
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.text = feature
        p = tf.paragraphs[0]
        p.font.size = Pt(14)
        p.font.color.rgb = DARK_COLOR
        p.alignment = PP_ALIGN.CENTER
    
    return slide

def create_presentation():
    """Create the complete CareVance HRMS presentation"""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    # ========== TITLE SLIDE ==========
    add_title_slide(
        prs,
        "CareVance HRMS",
        "Complete Workforce Operations Platform\nComprehensive Features & Modules Overview"
    )
    
    # ========== PLATFORM OVERVIEW ==========
    add_section_slide(prs, "Platform Overview")
    
    add_content_slide(prs, "About CareVance HRMS", [
        "Enterprise-grade Human Resource Management System designed for modern organizations",
        "Integrates time tracking, attendance monitoring, payroll management, and productivity analytics",
        "Built with cutting-edge technology stack: Laravel, React, Electron, and Expo",
        "Multi-platform solution: Web, Desktop, Mobile, and Browser Extension",
        "Suitable for businesses of all sizes - from startups to large enterprises",
        "Comprehensive employee lifecycle management from onboarding to exit",
        "Real-time monitoring and analytics for data-driven decision making",
        "Full compliance with Indian statutory regulations (PF, ESI, PT, TDS, LWF, Gratuity)"
    ])
    
    add_content_slide(prs, "Platform Components", [
        "Web Application (React 18 + TypeScript)",
        "  - Full-featured HRMS client with responsive design",
        "  - Intuitive dashboard and workspace interfaces",
        "Desktop Application (Electron 41)",
        "  - Time tracking with idle detection and auto-start",
        "  - Screenshot capture and activity monitoring",
        "Mobile Application (Expo/React Native)",
        "  - On-the-go attendance and leave management",
        "  - Push notifications and real-time updates",
        "Browser Extension (Chromium)",
        "  - Exact browser activity tracking",
        "  - Seamless integration with desktop tracker",
        "Backend API (Laravel 12 + PHP 8.2)",
        "  - RESTful API architecture",
        "  - Secure authentication and authorization"
    ])
    
    # ========== MODULE 1: TRACKING MANAGEMENT ==========
    add_section_slide(prs, "Module 1: Tracking Management")
    
    add_content_slide(prs, "Automatic Time Tracking", [
        "Desktop Timer with Intelligent Auto-Start",
        "  - Automatically starts at configured office hours",
        "  - Auto-launch on OS boot for seamless tracking",
        "  - Manual start/stop with single click",
        "Idle Time Detection & Management",
        "  - Configurable idle threshold (default: 3 minutes to track, 5 minutes to auto-stop)",
        "  - Lock screen detection - tracks time during screen lock",
        "  - Categorizes productive vs idle time automatically",
        "Real-time Activity Monitoring",
        "  - Tracks active applications and websites in real-time",
        "  - Desktop and browser activity separation",
        "  - Detailed activity logs with timestamps",
        "Timeline Visualization",
        "  - Visual timeline of daily work patterns",
        "  - Color-coded activity categories",
        "  - Quick overview of productive vs unproductive time"
    ])
    
    add_content_slide(prs, "Screenshots & Activity Capture", [
        "Automated Screenshot Capture",
        "  - Configurable capture intervals (default: every 10 minutes)",
        "  - Multiple monitor support for comprehensive tracking",
        "  - Blur/privacy mode for sensitive information",
        "Screenshot Management",
        "  - Bulk screenshot management with filters",
        "  - Detail view with zoom and navigation",
        "  - Secure signed URL access for privacy",
        "  - Screenshot deletion and retention policies",
        "Activity Sessions",
        "  - Tracks continuous work sessions",
        "  - Break time detection and categorization",
        "  - Session summaries with productivity metrics",
        "Offline Tracking",
        "  - Continues tracking without internet connection",
        "  - Automatic sync when connection restored",
        "  - No data loss during network interruptions"
    ])
    
    # ========== MODULE 2: PROJECT & TASK MANAGEMENT ==========
    add_section_slide(prs, "Module 2: Project & Task Management")
    
    add_content_slide(prs, "Project Management", [
        "Project Creation & Configuration",
        "  - Create projects with detailed descriptions and timelines",
        "  - Assign project managers and team members",
        "  - Set project budgets and hourly rates",
        "  - Define project priorities and categories",
        "Project Tracking & Monitoring",
        "  - Track time spent on each project",
        "  - Monitor project progress and milestones",
        "  - View project-specific reports and statistics",
        "  - Budget utilization and cost tracking",
        "Project Templates",
        "  - Save project configurations as templates",
        "  - Quick project setup using templates",
        "  - Standardized project workflows",
        "Project Archival",
        "  - Archive completed projects",
        "  - Historical project data retention",
        "  - Search and filter archived projects"
    ])
    
    add_content_slide(prs, "Task Management", [
        "Task Creation & Assignment",
        "  - Create tasks with titles, descriptions, and priorities",
        "  - Assign tasks to individual employees or teams",
        "  - Set due dates and reminders",
        "  - Attach files and documents to tasks",
        "Task Status & Workflow",
        "  - Multiple status options: To Do, In Progress, Review, Done",
        "  - Drag-and-drop status updates",
        "  - Task dependencies and blocking relationships",
        "  - Recurring task setup for regular activities",
        "Task Labels & Categorization",
        "  - Create custom labels for task organization",
        "  - Color-coded labels for visual identification",
        "  - Filter and sort tasks by labels",
        "Task Collaboration",
        "  - Task comments and discussions",
        "  - @mentions for team notifications",
        "  - Task activity history and audit trail",
        "  - Email notifications for task updates",
        "Time Tracking Integration",
        "  - Log time directly against specific tasks",
        "  - View time spent per task and user",
        "  - Task time reports and analytics"
    ])
    
    # ========== MODULE 3: ATTENDANCE MANAGEMENT ==========
    add_section_slide(prs, "Module 3: Attendance Management")
    
    add_content_slide(prs, "Check-in/Check-out System", [
        "Multiple Check-in Methods",
        "  - Web-based check-in from any browser",
        "  - Mobile app check-in with GPS verification",
        "  - Selfie-based attendance with face recognition",
        "  - Desktop timer auto-check-in option",
        "Today's Status Dashboard",
        "  - Real-time view of who's checked in/out",
        "  - Current work duration display",
        "  - Break status and duration tracking",
        "  - Late arrival and early departure alerts",
        "Calendar View",
        "  - Monthly attendance calendar with color coding",
        "  - Quick view of present, absent, and leave days",
        "  - Holiday and weekend markers",
        "  - Click to view daily details",
        "Attendance Summary Reports",
        "  - Monthly attendance statistics per employee",
        "  - Department-wise attendance summaries",
        "  - Late coming and early going reports",
        "  - Export reports to Excel/PDF"
    ])
    
    add_content_slide(prs, "Advanced Attendance Features", [
        "Holiday Management",
        "  - Company-wide holiday calendar",
        "  - Location-specific holidays",
        "  - Optional and mandatory holidays",
        "  - Holiday carryover policies",
        "Attendance Time Edit Requests",
        "  - Employees can request corrections to their attendance",
        "  - Manager approval workflow for edits",
        "  - Automatic overtime calculation for approved edits",
        "  - Audit trail of all changes",
        "Selfie Attendance with Geofencing",
        "  - Capture selfie during check-in/out",
        "  - GPS location verification",
        "  - Geofence zone restrictions",
        "  - Map view of check-in locations",
        "Shift Management",
        "  - Multiple shift configurations (Day, Night, Flexible)",
        "  - Shift rotation schedules",
        "  - Shift-wise attendance tracking",
        "  - Shift allowance calculations"
    ])
    
    # ========== MODULE 4: LEAVE MANAGEMENT ==========
    add_section_slide(prs, "Module 4: Leave Management")
    
    add_content_slide(prs, "Leave Request System", [
        "Leave Types Configuration",
        "  - Paid Time Off (PTO) / Casual Leave",
        "  - Sick Leave with medical certificate option",
        "  - Earned Leave / Privilege Leave",
        "  - Maternity and Paternity Leave",
        "  - Bereavement Leave",
        "  - Custom leave types as per company policy",
        "Leave Request Workflow",
        "  - Employee self-service leave application",
        "  - Date range selection with half-day option",
        "  - Reason and attachment upload",
        "  - Multi-level approval workflow (Manager → HR → Admin)",
        "Leave Balances & Entitlements",
        "  - Real-time leave balance display",
        "  - Annual leave entitlements by grade/tenure",
        "  - Leave accrual calculations",
        "  - Negative leave balance tracking",
        "Leave Calendar",
        "  - Visual calendar showing approved leaves",
        "  - Team leave view for planning",
        "  - Department-wise leave summary"
    ])
    
    add_content_slide(prs, "Advanced Leave Features", [
        "Leave Approval Workflow",
        "  - Email and in-app notifications for approvals",
        "  - Bulk approval option for managers",
        "  - Approval delegation during absence",
        "  - Comments and rejection reasons",
        "Leave Revocation",
        "  - Employees can request to cancel approved leaves",
        "  - Manager approval required for revocation",
        "  - Automatic balance restoration upon approval",
        "Compensatory Off (Comp-off)",
        "  - Track extra hours worked beyond regular schedule",
        "  - Automatic comp-off entitlement calculation",
        "  - Comp-off application and approval workflow",
        "  - Expiry tracking for unused comp-offs",
        "Leave Reports & Analytics",
        "  - Monthly/quarterly/annual leave reports",
        "  - Leave trend analysis",
        "  - Department-wise leave statistics",
        "  - Unused leave liability reports"
    ])
    
    # ========== MODULE 5: TEAM MANAGEMENT ==========
    add_section_slide(prs, "Module 5: Team Management")
    
    add_content_slide(prs, "Employee Management", [
        "Employee 360° Profiles",
        "  - Personal information and contact details",
        "  - Work information (department, designation, grade)",
        "  - Government ID documents (Aadhaar, PAN, Passport)",
        "  - Bank account details for salary processing",
        "  - Emergency contacts and dependents",
        "Organization Structure",
        "  - Department and sub-department hierarchy",
        "  - Designation and grade management",
        "  - Reporting manager assignment",
        "  - Team hierarchy visualization (Organization Tree)",
        "Employee Lifecycle",
        "  - New hires tracking and onboarding",
        "  - Employee status management (Active, Inactive, On Notice)",
        "  - Resignation workflow and exit processing",
        "  - Full & Final settlement management",
        "Groups Assignment",
        "  - Create functional and project groups",
        "  - Assign employees to multiple groups",
        "  - Group-based permissions and access",
        "  - Group-specific reports and analytics"
    ])
    
    # ========== MODULE 6: PAYROLL & FINANCE ==========
    add_section_slide(prs, "Module 6: Payroll & Finance")
    
    add_content_slide(prs, "Payroll Management Overview", [
        "Payroll Run Management",
        "  - Monthly payroll processing with multi-stage workflow",
        "  - Status tracking: Draft → Locked → Approved → Released → Paid",
        "  - Payroll preview before final processing",
        "  - Rollback and correction capabilities",
        "Salary Components & Structure",
        "  - Earnings: Basic, HRA, DA, CCA, Conveyance, Special Allowance",
        "  - Deductions: PF, ESI, PT, TDS, LWF, NPS, VPF",
        "  - Custom salary components with formulas",
        "  - Flexible salary templates by department/grade",
        "Pay Groups",
        "  - Multiple payroll groups with different pay schedules",
        "  - Monthly, bi-weekly, and weekly pay cycles",
        "  - Group-specific salary structures",
        "  - Department-wise payroll processing",
        "Payroll Automation",
        "  - Automatic payroll calculation based on attendance",
        "  - Integration with time tracking for accurate payments",
        "  - Automated statutory compliance calculations",
        "  - Scheduled payroll runs with notifications"
    ])
    
    add_content_slide(prs, "Statutory Compliance (India)", [
        "Provident Fund (PF)",
        "  - Employee and employer contribution calculations (12% each)",
        "  - EPS and EDLI component separation",
        "  - PF ECR (Electronic Challan cum Return) generation",
        "  - UAN management and tracking",
        "Employees' State Insurance (ESI)",
        "  - Automatic ESI applicability based on salary threshold",
        "  - Employee (0.75%) and employer (3.25%) contribution",
        "  - ESI Challan generation and filing",
        "  - IP number tracking",
        "Professional Tax (PT)",
        "  - State-wise PT slab configuration",
        "  - Monthly PT calculations based on gross salary",
        "  - PT challan generation for all applicable states",
        "Tax Deducted at Source (TDS)",
        "  - Monthly TDS calculations based on projected annual income",
        "  - New and Old tax regime support",
        "  - Form 16 generation at year-end",
        "  - Quarterly TDS return filing (24Q)"
    ])
    
    add_content_slide(prs, "Advanced Payroll Features", [
        "Tax Management",
        "  - Tax declaration submission by employees (80C, 80D, HRA, etc.)",
        "  - Tax proof upload and verification",
        "  - Tax simulator for regime comparison",
        "  - Tax-optimized salary structure suggestions",
        "Arrears & Backdated Payments",
        "  - Arrear calculations for salary revisions",
        "  - Backdated increment processing",
        "  - Recovery calculations for overpayments",
        "Leave Encashment",
        "  - Automatic leave encashment calculation",
        "  - Tax applicability on encashment amounts",
        "  - Encashment payout with regular salary",
        "Loans & Advances",
        "  - Employee loan management with EMI setup",
        "  - Salary advance processing",
        "  - Automatic EMI deduction from payroll",
        "  - Loan balance tracking and statements"
    ])
    
    add_content_slide(prs, "Payroll Additional Features", [
        "Reimbursements",
        "  - Expense claim submission with receipts",
        "  - Category-wise reimbursement limits",
        "  - Manager approval workflow",
        "  - Tax-free and taxable reimbursement separation",
        "Full & Final Settlement",
        "  - Automatic FnF calculation for resigned employees",
        "  - Notice period pay and recovery",
        "  - Leave encashment in settlement",
        "  - Gratuity calculation for eligible employees",
        "Flexible Benefits Program (FBP)",
        "  - Cafeteria-style benefits selection",
        "  - Component-wise allocation (LTA, Medical, etc.)",
        "  - FBP claims processing and validation",
        "  - Tax savings through optimal FBP structuring",
        "Variable Pay & Bonuses",
        "  - Performance-based variable pay (QVP, AVP)",
        "  - Annual bonus calculations",
        "  - Prorated bonus for new/exiting employees",
        "  - Retention and referral bonus management"
    ])
    
    add_content_slide(prs, "Payslip & Disbursement", [
        "Payslip Generation",
        "  - Detailed payslips with earnings and deductions breakdown",
        "  - YTD (Year-to-Date) summary in each payslip",
        "  - Multiple payslip templates",
        "  - Digital signature support",
        "Payslip Delivery",
        "  - Email payslip delivery with password protection",
        "  - Employee self-service payslip download",
        "  - Mobile app payslip access",
        "  - Payslip comparison between months",
        "Bank Integration",
        "  - Bank account validation and verification",
        "  - Direct bank transfer file generation",
        "  - Multiple bank support",
        "  - Razorpay integration for instant disbursement",
        "Payment Processing",
        "  - Bulk payment processing",
        "  - Payment status tracking",
        "  - Failed payment handling and retry",
        "  - Payment reconciliation with bank statements"
    ])
    
    # ========== MODULE 7: HRMS CORE ==========
    add_section_slide(prs, "Module 7: HRMS Core")
    
    add_content_slide(prs, "Employee Lifecycle Management", [
        "Onboarding",
        "  - Digital onboarding checklist and workflow",
        "  - Document collection and verification",
        "  - Welcome kit and asset assignment",
        "  - Buddy/mentor assignment",
        "  - Pre-boarding engagement before joining",
        "Employee Documents",
        "  - Centralized document repository",
        "  - Document categories: Personal, Education, Experience, etc.",
        "  - Expiry tracking for passports, visas, contracts",
        "  - Document version control",
        "Profile Management",
        "  - Employee self-service profile updates",
        "  - Approval workflow for critical changes",
        "  - Profile completeness tracking",
        "  - Professional social links integration",
        "Organization Structure",
        "  - Multi-level department hierarchy",
        "  - Cost center and profit center mapping",
        "  - Location-wise employee grouping",
        "  - Organization chart visualization"
    ])
    
    add_content_slide(prs, "Resignation & Exit Management", [
        "Resignation Workflow",
        "  - Employee resignation submission with notice period",
        "  - Manager acceptance and counter-offer option",
        "  - HR review and approval process",
        "  - Notice period waiver and adjustment",
        "Exit Interview",
        "  - Configurable exit interview questionnaire",
        "  - Interview scheduling and feedback capture",
        "  - Exit reason categorization and analytics",
        "  - Retention suggestions based on feedback",
        "Clearance Process",
        "  - Department-wise clearance checklist",
        "  - Asset return tracking",
        "  - IT account deactivation checklist",
        "  - Finance clearance for dues",
        "Full & Final Settlement",
        "  - Automatic settlement calculation",
        "  - Notice period pay/recovery",
        "  - Leave encashment and gratuity",
        "  - Settlement payout schedule",
        "Resignation Analytics",
        "  - Attrition rate calculation and trends",
        "  - Department-wise resignation analysis",
        "  - Exit reason trends",
        "  - Predictive attrition alerts"
    ])
    
    # ========== MODULE 8: COMMUNICATION ==========
    add_section_slide(prs, "Module 8: Communication & Engagement")
    
    add_content_slide(prs, "Chat & Messaging", [
        "Direct Messaging",
        "  - One-on-one private conversations",
        "  - Real-time message delivery with typing indicators",
        "  - Message edit and delete functionality",
        "  - File attachments with preview",
        "  - Message read receipts",
        "Group Chats",
        "  - Create groups for teams/projects",
        "  - Group admin and member management",
        "  - Group announcements and pinned messages",
        "  - Group message reactions (emoji)",
        "  - @mentions for specific members",
        "Notifications",
        "  - Desktop push notifications with sound",
        "  - Mobile push notifications via Expo",
        "  - Unread message badges and counters",
        "  - Notification preferences and muting",
        "Chat Features",
        "  - Search chat history",
        "  - Share screenshots and files",
        "  - Link previews",
        "  - Voice messages (future enhancement)"
    ])
    
    add_content_slide(prs, "Announcements & Engagement", [
        "Company Announcements",
        "  - Organization-wide announcement publishing",
        "  - Rich text editor with formatting",
        "  - Attachment support for documents/videos",
        "  - Priority levels (Normal, Important, Urgent)",
        "  - Read receipts and acknowledgment tracking",
        "Public Press & News",
        "  - External press release publishing",
        "  - Company news and achievements sharing",
        "  - Media gallery for events and celebrations",
        "  - Newsletter generation and distribution",
        "Polls & Surveys",
        "  - Quick polls for decision making",
        "  - Employee satisfaction surveys",
        "  - Anonymous feedback collection",
        "  - Survey result analytics and reporting",
        "Engagement Dashboard",
        "  - Recent activity feed",
        "  - Upcoming events and birthdays",
        "  - Work anniversaries celebration",
        "  - Recognition and kudos system"
    ])
    
    # ========== MODULE 9: PRODUCTIVITY & MONITORING ==========
    add_section_slide(prs, "Module 9: Productivity & Monitoring")
    
    add_content_slide(prs, "Productivity Management", [
        "Productivity Rules Engine",
        "  - Classify applications and websites by productivity",
        "  - Categories: Productive, Neutral, Unproductive",
        "  - Custom rules by department/role",
        "  - URL-level classification for granular control",
        "Productivity Scoring",
        "  - Individual employee productivity ratings",
        "  - Daily, weekly, and monthly productivity scores",
        "  - Trend analysis and improvement suggestions",
        "  - Gamification with productivity leaderboards",
        "Web & Application Usage",
        "  - Detailed web usage reports",
        "  - Most visited websites tracking",
        "  - Application usage time breakdown",
        "  - Category-wise usage statistics",
        "Idle Time Analysis",
        "  - Idle time tracking and categorization",
        "  - Idle patterns and behavior analysis",
        "  - Excessive idle time alerts",
        "  - Productive vs non-productive time ratio"
    ])
    
    add_content_slide(prs, "Monitoring & Reports", [
        "Live Monitoring Dashboard",
        "  - Real-time view of online employees",
        "  - Current activity status (Active, Idle, Offline)",
        "  - Screenshot thumbnails for quick review",
        "  - Work hours progress tracking",
        "Timeline Reports",
        "  - Visual timeline of employee workday",
        "  - Activity switches and context changes",
        "  - Break and idle time markers",
        "  - Export timeline as PDF/image",
        "Custom Reports",
        "  - Pre-built report templates",
        "  - Custom report builder with filters",
        "  - Scheduled report generation and email",
        "  - Export to Excel, PDF, and CSV formats",
        "Analytics & Insights",
        "  - Team productivity comparisons",
        "  - Department-wise performance metrics",
        "  - Time utilization analysis",
        "  - Attendance and punctuality trends"
    ])
    
    # ========== MODULE 10: BREAK MANAGEMENT ==========
    add_section_slide(prs, "Module 10: Break Management")
    
    add_content_slide(prs, "Break Tracking System", [
        "Break Types",
        "  - Lunch break tracking",
        "  - Tea/Coffee break tracking",
        "  - Short break management",
        "  - Custom break types as per policy",
        "Break Initiation",
        "  - One-click break start/stop",
        "  - Automatic break reminders",
        "  - Maximum break duration alerts",
        "  - Break approval for extended breaks",
        "Break Reports",
        "  - Daily break duration summary",
        "  - Average break time analytics",
        "  - Excessive break alerts",
        "  - Break pattern analysis",
        "Break Policies",
        "  - Configurable break limits per shift",
        "  - Paid vs unpaid break categorization",
        "  - Break deduction in payroll calculation",
        "  - Grace period configuration"
    ])
    
    # ========== MODULE 11: OVERTIME MANAGEMENT ==========
    add_section_slide(prs, "Module 11: Overtime Management")
    
    add_content_slide(prs, "Overtime Tracking & Calculation", [
        "Overtime Rules Configuration",
        "  - Daily/weekly overtime thresholds",
        "  - Overtime calculation methods (Daily, Weekly, Monthly)",
        "  - Different OT rates (1.5x, 2x, 3x) by day type",
        "  - Comp-off vs payment options",
        "Automatic OT Calculation",
        "  - Automatic detection of overtime hours",
        "  - Integration with attendance and time tracking",
        "  - OT approval workflow by manager",
        "  - OT payment in payroll or comp-off credit",
        "Overtime Reports",
        "  - Monthly OT summary by employee",
        "  - Department-wise OT analysis",
        "  - OT cost analysis and budgeting",
        "  - OT trends and patterns",
        "OT Policies",
        "  - Pre-approval requirement for OT",
        "  - Maximum OT limits per day/week",
        "  - Mandatory rest period after OT",
        "  - Weekend and holiday OT handling"
    ])
    
    # ========== MODULE 12: APPROVAL WORKFLOW ==========
    add_section_slide(prs, "Module 12: Approval Management")
    
    add_content_slide(prs, "Multi-level Approval System", [
        "Centralized Approval Inbox",
        "  - Single dashboard for all pending approvals",
        "  - Filter by type (Leave, OT, Expense, etc.)",
        "  - Bulk approval/rejection option",
        "  - Approval history and audit trail",
        "Configurable Approval Chains",
        "  - Define approval hierarchy by request type",
        "  - Multi-level approvals (Manager → HR → Finance → Admin)",
        "  - Parallel approval for different departments",
        "  - Approval delegation during absence",
        "Approval Types",
        "  - Leave requests (Casual, Sick, Earned, etc.)",
        "  - Attendance time edit requests",
        "  - Expense and reimbursement claims",
        "  - Loan and advance requests",
        "  - Resignation and exit clearances",
        "Notifications & Reminders",
        "  - Instant email and in-app notifications",
        "  - Daily pending approval reminders",
        "  - Escalation for overdue approvals",
        "  - Approval completion confirmation"
    ])
    
    # ========== MODULE 13: INTEGRATIONS & API ==========
    add_section_slide(prs, "Module 13: Integrations & API")
    
    add_content_slide(prs, "API & Integration Capabilities", [
        "Open API Access",
        "  - RESTful API with comprehensive endpoints",
        "  - API documentation with Swagger/OpenAPI",
        "  - API key-based authentication",
        "  - Rate limiting and usage analytics",
        "Third-Party Integrations",
        "  - Google OAuth for single sign-on",
        "  - Razorpay for payment processing",
        "  - Stripe for subscription billing",
        "  - Email providers (SMTP, SendGrid, AWS SES)",
        "AI Integration",
        "  - AI-powered productivity insights",
        "  - Smart attendance predictions",
        "  - Automated anomaly detection",
        "  - Chatbot integration for HR queries",
        "Mobile App Integration",
        "  - Native iOS and Android apps",
        "  - Push notification service",
        "  - Biometric authentication support",
        "  - Offline data sync capability",
        "Custom Integrations",
        "  - Webhook support for real-time events",
        "  - Custom connector development",
        "  - Data import/export APIs",
        "  - Single Sign-On (SSO) support"
    ])
    
    # ========== MODULE 14: REPORTS & ANALYTICS ==========
    add_section_slide(prs, "Module 14: Reports & Analytics")
    
    add_content_slide(prs, "Comprehensive Reporting Suite", [
        "Dashboard Analytics",
        "  - Real-time dashboard with KPIs",
        "  - Attendance and punctuality metrics",
        "  - Productivity trends and comparisons",
        "  - Payroll summary and compliance status",
        "Report Types",
        "  - Time Reports: Daily, Weekly, Monthly timesheets",
        "  - Attendance Reports: Summary, detailed, and exception reports",
        "  - Leave Reports: Balance, consumption, and trends",
        "  - Payroll Reports: CTC, variance, YTD, and MIS",
        "  - Custom Reports: User-defined report builder",
        "Advanced Analytics",
        "  - Employee timeline and journey analysis",
        "  - Department comparison reports",
        "  - Cost center analysis",
        "  - ROI on workforce investments",
        "Export & Sharing",
        "  - Export to Excel, PDF, and CSV formats",
        "  - Scheduled report generation",
        "  - Email distribution lists",
        "  - Secure report sharing with access control"
    ])
    
    # ========== MODULE 15: SECURITY & COMPLIANCE ==========
    add_section_slide(prs, "Module 15: Security & Compliance")
    
    add_content_slide(prs, "Security Features", [
        "Authentication & Authorization",
        "  - Role-based access control (RBAC)",
        "  - Multi-factor authentication (MFA) support",
        "  - Session management and timeout",
        "  - Password policy enforcement",
        "Data Security",
        "  - AES-256 encryption for sensitive data",
        "  - SSL/TLS for data in transit",
        "  - Database encryption at rest",
        "  - Secure file storage with access controls",
        "Audit & Compliance",
        "  - Comprehensive audit logging",
        "  - User activity tracking",
        "  - Data access logs",
        "  - GDPR and privacy compliance features",
        "Infrastructure Security",
        "  - Regular security audits and penetration testing",
        "  - Automated vulnerability scanning",
        "  - DDoS protection",
        "  - Automatic security patches and updates",
        "Access Controls",
        "  - IP-based access restrictions",
        "  - Device registration and management",
        "  - Geo-location based access control",
        "  - Time-based access restrictions"
    ])
    
    # ========== MODULE 16: MOBILE APP ==========
    add_section_slide(prs, "Module 16: Mobile Application")
    
    add_content_slide(prs, "Mobile App Features", [
        "Attendance on Mobile",
        "  - GPS-based check-in/out",
        "  - Selfie attendance with face verification",
        "  - Check-in history and calendar view",
        "  - Offline check-in with auto-sync",
        "Leave Management",
        "  - Apply for leave on the go",
        "  - View leave balances and history",
        "  - Leave approval from mobile",
        "  - Leave calendar and team view",
        "Payslip & Payroll",
        "  - View and download payslips",
        "  - Tax declaration submission",
        "  - Reimbursement claims",
        "  - Loan and advance requests",
        "Communication",
        "  - Push notifications for approvals",
        "  - Chat and announcements",
        "  - Company news and updates",
        "  - Emergency contact access",
        "Employee Self-Service",
        "  - Profile management",
        "  - Document upload",
        "  - Shift schedule view",
        "  - Task and project updates"
    ])
    
    # ========== SUMMARY SLIDES ==========
    add_section_slide(prs, "Summary & Key Benefits")
    
    add_content_slide(prs, "Key Benefits", [
        "For Management:",
        "  • Complete visibility into workforce productivity and attendance",
        "  • Data-driven decision making with comprehensive analytics",
        "  • Reduced administrative overhead through automation",
        "  • Compliance assurance with statutory regulations",
        "",
        "For HR Teams:",
        "  • Streamlined payroll processing with 100% accuracy",
        "  • Simplified leave and attendance management",
        "  • Centralized employee data and documentation",
        "  • Automated workflows reduce manual tasks by 70%",
        "",
        "For Employees:",
        "  • Self-service portal for all HR needs",
        "  • Transparent leave balances and payslip access",
        "  • Fair and accurate time tracking",
        "  • Mobile app for on-the-go access",
        "",
        "For IT Teams:",
        "  • Secure, cloud-based solution with minimal maintenance",
        "  • Role-based access control and audit trails",
        "  • API-first architecture for easy integrations",
        "  • Regular updates and security patches"
    ])
    
    add_content_slide(prs, "Why Choose CareVance HRMS?", [
        "Comprehensive Solution:",
        "  • All-in-one platform covering entire employee lifecycle",
        "  • 16+ integrated modules with seamless data flow",
        "  • No need for multiple disjointed systems",
        "",
        "Technology Leadership:",
        "  • Modern tech stack: React, Laravel, Electron, Expo",
        "  • Multi-platform support: Web, Desktop, Mobile, Browser",
        "  • AI-powered insights and automation",
        "",
        "Scalability & Flexibility:",
        "  • Suitable for startups to enterprise organizations",
        "  • Modular architecture - use only what you need",
        "  • Customizable workflows and policies",
        "",
        "Security & Compliance:",
        "  • Enterprise-grade security with encryption",
        "  • Full compliance with Indian labor laws",
        "  • Regular security audits and updates",
        "",
        "Support & Service:",
        "  • 24/7 customer support",
        "  • Dedicated account management",
        "  • Regular training and onboarding assistance"
    ])
    
    # Final slide
    add_title_slide(
        prs,
        "Thank You",
        "CareVance HRMS - Transforming Workforce Management\n\nContact us for a personalized demo\nwww.carevance.com | support@carevance.com"
    )
    
    return prs

if __name__ == "__main__":
    print("Creating CareVance HRMS Presentation...")
    presentation = create_presentation()
    
    # Save the presentation
    output_path = r"D:\CareVance_Hrms_IDE\CareVance_HRMS_Features_Presentation.pptx"
    presentation.save(output_path)
    
    print(f"[SUCCESS] Presentation created successfully!")
    print(f"[FILE] File saved to: {output_path}")
    print(f"[COUNT] Total slides: {len(presentation.slides)}")
